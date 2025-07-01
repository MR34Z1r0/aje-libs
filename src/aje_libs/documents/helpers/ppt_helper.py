# src/aje_libs/bd/helpers/pptx_helper.py
from pptx import Presentation
from typing import List, Dict, Optional
import os
import tempfile
from pathlib import Path
from ...common.logger import custom_logger

logger = custom_logger(__name__)

class PPTXHelper:
    """Helper para procesar archivos PPTX y PPT"""
    
    def __init__(self):
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Verifica las dependencias disponibles"""
        self.has_win32 = False
        try:
            import win32com.client
            self.has_win32 = True
            logger.info("pywin32 disponible para conversión PPT")
        except ImportError:
            logger.warning("pywin32 no disponible. Solo se procesarán archivos PPTX")
    
    def _is_ppt_file(self, file_path: str) -> bool:
        """Determina si el archivo es PPT basándose en la extensión"""
        return Path(file_path).suffix.lower() == '.ppt'
    
    def _convert_ppt_to_pptx(self, ppt_path: str) -> str:
        """
        Convierte archivo PPT a PPTX usando COM automation (Windows)
        
        :param ppt_path: Ruta al archivo PPT
        :return: Ruta al archivo PPTX temporal
        """
        if not self.has_win32:
            raise RuntimeError("pywin32 requerido para procesar archivos PPT. Instala con: pip install pywin32")
        
        try:
            import win32com.client
            
            # Crear archivo temporal para PPTX
            temp_dir = tempfile.gettempdir()
            temp_pptx = os.path.join(temp_dir, f"temp_{os.path.basename(ppt_path)}.pptx")
            
            # Inicializar PowerPoint
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False
            
            try:
                # Abrir presentación PPT
                presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
                
                # Guardar como PPTX (formato 24 = ppSaveAsOpenXMLPresentation)
                presentation.SaveAs(os.path.abspath(temp_pptx), 24)
                presentation.Close()
                
                logger.info(f"PPT convertido a PPTX temporal: {temp_pptx}")
                return temp_pptx
                
            finally:
                powerpoint.Quit()
                
        except Exception as e:
            logger.error(f"Error convirtiendo PPT a PPTX: {e}")
            raise e
    
    def _cleanup_temp_file(self, file_path: str):
        """Elimina archivo temporal si existe"""
        try:
            if os.path.exists(file_path) and "temp_" in os.path.basename(file_path):
                os.remove(file_path)
                logger.debug(f"Archivo temporal eliminado: {file_path}")
        except Exception as e:
            logger.warning(f"No se pudo eliminar archivo temporal {file_path}: {e}")
    
    def extract_text(self, file_path: str) -> str:
        """
        Extrae texto de una presentación PPTX o PPT.
        
        :param file_path: Ruta al archivo PPTX o PPT
        :return: Texto extraído
        """
        temp_file = None
        try:
            # Determinar si necesita conversión
            if self._is_ppt_file(file_path):
                file_path = self._convert_ppt_to_pptx(file_path)
                temp_file = file_path
            
            prs = Presentation(file_path)
            all_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:  # Solo agregar si hay contenido
                    all_text.append(f"Slide {i+1}: {' '.join(slide_text)}")
            
            extracted_text = "\n\n".join(all_text)
            logger.info(f"Texto extraído exitosamente: {len(all_text)} diapositivas")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Error extrayendo texto: {e}")
            raise e
        finally:
            if temp_file:
                self._cleanup_temp_file(temp_file)
    
    def extract_slides(self, file_path: str) -> List[Dict[str, str]]:
        """
        Extrae texto por diapositiva de PPTX o PPT.
        
        :param file_path: Ruta al archivo PPTX o PPT
        :return: Lista de diapositivas con su contenido
        """
        temp_file = None
        try:
            # Determinar si necesita conversión
            if self._is_ppt_file(file_path):
                file_path = self._convert_ppt_to_pptx(file_path)
                temp_file = file_path
            
            prs = Presentation(file_path)
            slides = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = {
                    'slide_number': i + 1,
                    'title': '',
                    'content': []
                }
                
                shapes_with_text = [shape for shape in slide.shapes 
                                  if hasattr(shape, 'text') and shape.text.strip()]
                
                if shapes_with_text:
                    # El primer shape con texto suele ser el título
                    slide_content['title'] = shapes_with_text[0].text.strip()
                    
                    # El resto son contenido
                    for shape in shapes_with_text[1:]:
                        if shape.text.strip():
                            slide_content['content'].append(shape.text.strip())
                
                slides.append(slide_content)
            
            logger.info(f"Extraídas {len(slides)} diapositivas")
            return slides
            
        except Exception as e:
            logger.error(f"Error extrayendo diapositivas: {e}")
            raise e
        finally:
            if temp_file:
                self._cleanup_temp_file(temp_file)
    
    def get_slide_layouts(self, file_path: str) -> List[Dict[str, any]]:
        """
        Obtiene información sobre los layouts de las diapositivas.
        
        :param file_path: Ruta al archivo PPTX o PPT
        :return: Lista de layouts utilizados
        """
        temp_file = None
        try:
            # Determinar si necesita conversión
            if self._is_ppt_file(file_path):
                file_path = self._convert_ppt_to_pptx(file_path)
                temp_file = file_path
            
            prs = Presentation(file_path)
            layouts = []
            
            for i, slide in enumerate(prs.slides):
                layout_info = {
                    'slide_number': i + 1,
                    'layout_name': slide.slide_layout.name,
                    'shapes_count': len(slide.shapes),
                    'text_shapes_count': len([s for s in slide.shapes 
                                            if hasattr(s, 'text') and s.text.strip()])
                }
                layouts.append(layout_info)
            
            logger.info(f"Información de layouts obtenida para {len(layouts)} diapositivas")
            return layouts
            
        except Exception as e:
            logger.error(f"Error obteniendo layouts: {e}")
            raise e
        finally:
            if temp_file:
                self._cleanup_temp_file(temp_file)
    
    def get_supported_formats(self) -> List[str]:
        """Retorna los formatos soportados"""
        formats = ['.pptx']
        if self.has_win32:
            formats.append('.ppt')
        return formats
    
    def is_supported_format(self, file_path: str) -> bool:
        """Verifica si el formato del archivo es soportado"""
        extension = Path(file_path).suffix.lower()
        return extension in self.get_supported_formats()